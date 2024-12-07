from pptx import Presentation
import pytesseract
from PIL import Image
import re
import camelot
import json


class TableExtractor:
    def __init__(self, config):
        self.config = config

    def extract_tables_from_pdf(self, pdf_path: str):
        # Извлекаем таблицы из PDF
        try:
            tables = camelot.read_pdf(pdf_path, pages='all',
                                      flavor='stream')  # Используем flavor='stream' для текстовых таблиц
            if tables:
                tables_data = [table.df.to_dict(orient="records") for table in tables]
                return json.dumps(tables_data, ensure_ascii=False)
            else:
                return "No tables found in this PDF"
        except Exception as e:
            return f"Error extracting tables from PDF: {str(e)}"

    def extract_tables_from_pptx(self, pptx_path: str):
        # Извлекаем таблицы из PowerPoint
        prs = Presentation(pptx_path)
        tables = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "table"):  # Проверяем на наличие таблицы в объекте
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    tables.append(table_data)
        if tables:
            return json.dumps(tables, ensure_ascii=False)
        else:
            return "No tables found in this PPTX file"

    def extract_tables_from_image(self, image_path: str):
        # Извлекаем таблицы из изображения с помощью OCR
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img)

        # Простой способ поиска таблиц (через пробелы или табуляцию)
        table_pattern = r'([^\n]+(?:\t[^\n]+)+)'  # Пример для табуляций
        rows = re.findall(table_pattern, text)

        if rows:
            table_data = [row.split("\t") for row in rows]
            return json.dumps(table_data, ensure_ascii=False)
        else:
            return "No tables found in this image"
