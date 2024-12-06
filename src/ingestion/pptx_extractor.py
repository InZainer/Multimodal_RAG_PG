# src/ingestion/pptx_extractor.py

from pptx import Presentation

class PPTXExtractor:
    @staticmethod
    def extract_text(pptx_path: str) -> str:
        prs = Presentation(pptx_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
